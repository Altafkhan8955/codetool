from django.shortcuts import render
from django.template import RequestContext
from app.models import Review,ClientQuery,Service
from django.views import View
from django.conf import settings

# Create your views here.
def index(request):
    #ip = get_client_ip(request)
    service = Service.objects.all()
    ip=request.META.get("REMOTE_ADDR",None)
    context = {
        "ip":ip,
        "service":service,
    }
    return render(request, "index.html",context)
def ClientReview(request):
    if request.method == "POST":
        name = request.POST.get("name")
        email = request.POST.get("email")
        comment = request.POST.get("comment")
        ip = request.META.get("REMOTE_ADDR",None)
        service = Service.objects.all()
        review = Review(name=name, email=email, comment=comment, ip=ip)
        review.save()
        msg = "Your review successfull submit"
        context = {
            "ip":ip,
            "service":service,
            "msg":msg,
        }
        return render(request,"index.html",context)
  
def Contact(request):
    if request.method == "POST":
        name = request.POST.get("name")
        email = request.POST.get("email")
        message = request.POST.get("message")
        ip = request.META.get("REMOTE_ADDR",None)
        service = Service.objects.all()
        review = Review.objects.all()
        query = ClientQuery(name=name, email=email, message=message, ip=ip)
        query.save()
        msg = "Your query successfull submit"
        context = {
            "ip":ip,
            "service":service,
            "review":review,
            "msg":msg,
        }
        return render(request,"contact.html",context)
    else:
         return render(request,"contact.html")
        
def custom_bad_request_view(request, exception):
    return render(request, '400.html', status=404)
def custom_unthorized_view(request, exception):
    return render(request,'401.html', status=401)
def custom_permission_denied_view(request, exception):
    return render(request, '403.html', status=403)
def custom_page_not_found_view(request, exception):
    return render(request, '404.html', status=404)
def custom_server_error_view(request):
    return render(request, '500.html', status=500)
def sitemap_view(request):
    return render(request, 'sitemap.xml', content_type='application/xml')
def sitemap0_view(request):
    return render(request, 'sitemap-0.xml', content_type='application/xml')
#########################image tool################################
from django.shortcuts import render, redirect, get_object_or_404
from django.http import FileResponse
from .forms import ImageForm,ImageResizeForm,ResizeForm,ConvertpngForm,ConvertjpgForm,ConvertzipForm,KeywordForm,URLForm
from app.models import ImageModel,ResizableImage,ResizedImage,ConvertpngModel,ConvertjpgModel
from PIL import Image, ImageFilter
import rembg
from io import BytesIO
from django.http import HttpResponse
import os
import zipfile
from zipfile import ZipFile
import mimetypes
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from PyPDF2 import PdfMerger,PdfReader
from fpdf import FPDF
from subprocess import PIPE, Popen
import subprocess
from pdf2docx import Converter
import fitz
import csv
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image, ImageDraw, ImageFont
import base64
from django.core.files.storage import FileSystemStorage
from django.core.files.uploadedfile import InMemoryUploadedFile
import numpy as np
import cv2
#remove image background
def remove_background(image_path):
    # Use rembg to remove the background
    with open(image_path, 'rb') as input_file:
        input_data = input_file.read()
    output_data = rembg.remove(input_data)
    # Convert the output data to a PIL Image
    processed_image = Image.open(BytesIO(output_data))
    return processed_image
def process_image(request):
    if request.method == 'POST':
        form = ImageForm(request.POST, request.FILES)
        if form.is_valid():
            image_model = form.save()
            # Call your background removal function
            processed_image = remove_background(image_model.image.path)
            # Serve the processed image for download
            response = HttpResponse(content_type='image/png')
            processed_image.save(response, format='PNG')
            response['Content-Disposition'] = f'attachment; filename={image_model.image.name}'
            return response
    else:
        form = ImageForm()
    return render(request, 'imagetool/background-remove.html', {'form': form})
#Resize by mb and kb
def resize_image(image, width, height):
    img = Image.open(image)
    img = img.resize((width, height)) #Image.ANTIALIAS
    output = BytesIO()
    img.save(output, format='JPEG')
    return output.getvalue()
def upload_and_resize(request):
    if request.method == 'POST':
        form = ImageResizeForm(request.POST, request.FILES)
        if form.is_valid():
            original_image = form.cleaned_data['image']
            # Save the original image
            instance = ResizableImage(original_image=original_image)
            instance.save()
            # Resize the image
            resized_image_data = resize_image(original_image, 300, 200)
            return render(request, 'imagetool/image-resize.html', {'instance': instance, 'resized_image_data': resized_image_data})
    else:
        form = ImageResizeForm()
    return render(request, 'imagetool/image-resize.html', {'form': form})
def download_resized_image(request, pk):
    image_instance = get_object_or_404(ResizableImage, pk=pk)
    resized_image_data = resize_image(image_instance.original_image, 300, 200)
    response = HttpResponse(resized_image_data, content_type='image/jpeg')
    response['Content-Disposition'] = f'attachment; filename="{image_instance.original_image}_resized.jpg"'
    return response
#resize image by height and width
def resizedimage(request):
    if request.method == 'POST':
        form = ResizeForm(request.POST, request.FILES)
        if form.is_valid():
            height = form.cleaned_data['height']
            width = form.cleaned_data['width']
            image = form.cleaned_data['image']
            # Resize the image
            img = Image.open(image)
            img_resized = img.resize((width, height))
            # Save the resized image
            resized_image_io = BytesIO()
            img_resized.save(resized_image_io, format='JPEG')
            img_name = os.path.splitext(image.name)[0] + f'_resized_{width}x{height}.jpg'
            resized_image = ResizedImage(image=img_name)
            resized_image.image.save(img_name, content=BytesIO(resized_image_io.getvalue()))
            resized_image.save()
            return redirect('download_resize', pk=resized_image.pk)
    else:
        form = ResizeForm()
    return render(request, 'imagetool/resize_image.html', {'form': form})
def download_resize(request, pk):
    resized_image = ResizedImage.objects.get(pk=pk)
    response = HttpResponse(resized_image.image.read(), content_type='image/jpeg')
    response['Content-Disposition'] = f'attachment; filename="{resized_image.image.name}"'
    return response
#convert image format jpg to png 
def convertpng_image(request, pk):
    image_model = get_object_or_404(ConvertpngModel, pk=pk)
    image = Image.open(image_model.image.path)
    # Convert image to PNG
    png_image = BytesIO()
    image.save(png_image, format='PNG')
    png_image.seek(0)
    # Prepare response with PNG content
    response = HttpResponse(png_image, content_type='image/png')
    response['Content-Disposition'] = f'attachment; filename={image_model.image.name.replace(" ", "_").lower()}.png'
    return response
def uploadpng_image(request):
    if request.method == 'POST':
        form = ConvertpngForm(request.POST, request.FILES)
        if form.is_valid():
            instance = form.save()
            return render(request, 'imagetool/convertpng.html', {'image_model': instance})
    else:
        form = ConvertpngForm()
    return render(request, 'imagetool/convertpng.html', {'form': form})
#convert image format jpg
def convertjpg_image(png_image):
    # Convert PNG to JPG
    img = Image.open(png_image.path)
    jpg_io = BytesIO()
    img.save(jpg_io, format='JPEG')
    jpg_image = Image.open(jpg_io)
    # Save the converted image
    jpg_path = os.path.splitext(png_image.path)[0] + '.jpg'
    jpg_image.save(jpg_path)
    return jpg_path

def downloadjpg_image(request, image_path):
    # Provide the image for download
    with open(image_path, 'rb') as f:
        response = HttpResponse(f.read(), content_type='image/jpeg')
        response['Content-Disposition'] = 'inline; filename=' + os.path.basename(image_path)
        return response
def uploadjpg_image(request):
    if request.method == 'POST':
        form = ConvertjpgForm(request.POST, request.FILES)
        if form.is_valid():
            image = form.save()
            jpg_path = convertjpg_image(image.image)
            return render(request, 'imagetool/convertjpg.html', {'jpg_path': jpg_path})
    else:
        form = ConvertjpgForm()
    return render(request, 'imagetool/convertjpg.html', {'form': form})
#convert image to zip format
def Convertzip_and_download(request):
    if request.method == 'POST':
        form = ConvertzipForm(request.POST, request.FILES)
        if form.is_valid():
            images = form.cleaned_data['images']
           # images = request.FILES.getlist('images')
            # Create a temporary directory
            temp_dir = 'temp_images'
            os.makedirs(temp_dir, exist_ok=True)
            # Convert images and save to the temporary directory
            converted_images = []
            for image in images:
                img = Image.open(image)
                converted_img = os.path.join(temp_dir, f"{image.name.split('.')[0]}.png")
                img.save(converted_img, 'PNG')
                converted_images.append(converted_img)
            # Create a zip file
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
                for converted_img in converted_images:
                    zip_file.write(converted_img, os.path.basename(converted_img))
            # Cleanup temporary directory
            for converted_img in converted_images:
                os.remove(converted_img)
            os.rmdir(temp_dir)
            # Prepare response
            zip_buffer.seek(0)
            response = HttpResponse(zip_buffer, content_type='application/zip')
            response['Content-Disposition'] = 'attachment; filename=converted_images.zip'
            return response
    else:
        form = ConvertzipForm()
    return render(request, 'imagetool/convertzip.html', {'form': form})
#convert etract zip file
def Extractzip(request):
    file_list = request.session.get('file_list', [])
    return render(request, 'imagetool/extractzip.html', {'file_list': file_list})
def extract_zip(request):
    if request.method == 'POST':
        zip_file = request.FILES.get('zip_file')
        if zip_file and zip_file.name.endswith('.zip'):
            # Extract image files to a temporary directory
            temp_dir = 'temp_extracted_images'
            os.makedirs(temp_dir, exist_ok=True)
            with ZipFile(zip_file, 'r') as zip_ref:
                for file_name in zip_ref.namelist():
                    if file_name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                        zip_ref.extract(file_name, temp_dir)
            # Save the list of extracted image files to the session
            file_list = [name for name in os.listdir(temp_dir) if name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
            request.session['file_list'] = file_list
    return redirect('Extractzip')
def downloadextractzip_image(request, filename):
    temp_dir = 'temp_extracted_images'
    file_path = os.path.join(temp_dir, filename)
    # Serve the image file for download
    content_type, _ = mimetypes.guess_type(file_path)
    response = FileResponse(open(file_path, 'rb'), content_type=content_type)
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response
#convert pdf to image
def convert_pdf_to_image(pdf_path, output_path):
    pdf_document = fitz.open(pdf_path)
    page = pdf_document[0]  # Assuming you want to convert the first page
    image_list = page.get_images(full=True)
    image_index = 0  # You can change this to convert a different image if there are multiple images on the page
    image = page.get_pixmap()
    image_path = os.path.join(output_path, 'output.png')
    image.save(image_path)
    return image_path
def pdf_to_image(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        pdf_path = os.path.join('media', 'input.pdf')
        output_path = os.path.join('media', 'output')
        # Create the output directory if it doesn't exist
        os.makedirs(output_path, exist_ok=True)
        with open(pdf_path, 'wb') as destination:
            for chunk in pdf_file.chunks():
                destination.write(chunk)
        image_path = convert_pdf_to_image(pdf_path, output_path)
        os.remove(pdf_path)
        # Check if the image file exists before attempting to open it
        if os.path.exists(image_path):
            with open(image_path, 'rb') as image_file:
                response = HttpResponse(image_file.read(), content_type='image/png')
                response['Content-Disposition'] = 'attachment; filename=output.png'
                image_file.close()
                os.remove(image_path)
            return response
        else:
            return HttpResponse("Conversion failed. The image file is not available.", status=500)
    return render(request, 'imagetool/pdf_to_image.html')
#image compressor
def ImageCompressor(request):
    context = {}
    if request.method == 'POST' and request.FILES.get('image'):
        image_file = request.FILES['image']
        # Open image using PIL
        img = Image.open(image_file)
        # Convert to RGB if PNG with transparency (important)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        # Compress image (adjust quality here)
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG', optimize=True, quality=85)  
        buffer.seek(0)
        # Convert to base64 for preview
        import base64
        img_base64 = base64.b64encode(buffer.getvalue()).decode()
        context['image'] = img_base64
        # Store compressed image in session (temporary)
        request.session['compressed_image'] = img_base64
    return render(request, 'imagetool/compressor.html', context)
def DownloadCompressor(request):
    img_data = request.session.get('compressed_image')
    if not img_data:
        return HttpResponse("No image found")
    image_bytes = base64.b64decode(img_data)
    response = HttpResponse(image_bytes, content_type='image/jpeg')
    response['Content-Disposition'] = 'attachment; filename="compressed.jpg"'
    return response
##############################image watermark adder############################
def safe_int(value, default=0):
    try:
        return int(float(value))
    except (TypeError, ValueError):
        return default
def ImageWatermark(request):
    image_data = None
    if request.method == 'POST':
        image_file = request.FILES.get('image')
        watermark_text = request.POST.get('watermark')
        # Position from frontend
        pos_x = safe_int(request.POST.get('pos_x', 0))
        pos_y = safe_int(request.POST.get('pos_y', 0))
        if image_file:
            image = Image.open(image_file).convert("RGBA")
            txt_layer = Image.new("RGBA", image.size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(txt_layer)
            try:
                font = ImageFont.truetype("arial.ttf", 70)
            except:
                font = ImageFont.load_default()
            draw.text((pos_x, pos_y), watermark_text,
                      fill=(255, 255, 255, 150), font=font)
            watermarked = Image.alpha_composite(image, txt_layer)
            final_image = watermarked.convert("RGB")
            buffer = io.BytesIO()
            final_image.save(buffer, format='JPEG')
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            image_data = image_base64
    return render(request, 'imagetool/watermark.html', {'image_data': image_data})
# image watermark remover
def RemoveWatermark(image):
    # Convert to OpenCV format
    img = np.array(image)
    img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)
    # Convert to grayscale
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    # Threshold to detect watermark (works for light text)
    _, thresh = cv2.threshold(gray, 200, 255, cv2.THRESH_BINARY)
    # Inpaint to remove watermark
    result = cv2.inpaint(img, thresh, 3, cv2.INPAINT_TELEA)
    return result
def ImageWatermarkRemover(request):
    context = {}
    if request.method == 'POST' and request.FILES.get('image'):
        uploaded_file = request.FILES['image']
        # Open image using PIL
        image = Image.open(uploaded_file).convert('RGB')
        # Remove watermark
        processed_img = RemoveWatermark(image)
        # Convert back to RGB
        processed_img = cv2.cvtColor(processed_img, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(processed_img)
        # Convert to base64 (NO saving)
        buffer = BytesIO()
        pil_img.save(buffer, format="PNG")
        img_str = base64.b64encode(buffer.getvalue()).decode()
        context['image'] = img_str
    return render(request, 'imagetool/watermark-remover.html', context)
##############################image face blur###################################
def ImageFaceBlur(request):
    blurred_image = None
    if request.method == 'POST' and request.FILES.get('image'):
        uploaded_image = request.FILES['image']
        # Convert uploaded image to numpy array
        file_bytes = np.asarray(bytearray(uploaded_image.read()), dtype=np.uint8)
        img = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)
        # Load OpenCV Face Detector
        face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        faces = face_cascade.detectMultiScale(gray,scaleFactor=1.1,minNeighbors=5,minSize=(40, 40))
        # Advanced Blur
        for (x, y, w, h) in faces:
            face = img[y:y+h, x:x+w]
            # Strong Gaussian Blur
            blurred_face = cv2.GaussianBlur(face, (99, 99), 30)
            img[y:y+h, x:x+w] = blurred_face
        # Convert image to base64
        _, buffer = cv2.imencode('.jpg', img)
        image_base64 = base64.b64encode(buffer).decode('utf-8')
        blurred_image = image_base64
    return render(request, 'imagetool/image_blur.html', {'blurred_image': blurred_image})
##############################Image background blur#############################
def BackgroundBlur(request):
    blurred_image = None
    if request.method == 'POST' and request.FILES.get('image'):
        uploaded_image = request.FILES['image']
        # Read image bytes
        input_bytes = uploaded_image.read()
        # Remove background using AI
        output = rembg.remove(input_bytes)
        # Convert original image to OpenCV
        original_np = np.frombuffer(input_bytes, np.uint8)
        original_img = cv2.imdecode(original_np, cv2.IMREAD_COLOR)
        # Convert removed background image
        removed_bg = Image.open(io.BytesIO(output)).convert("RGBA")
        # Convert PIL to numpy
        removed_np = np.array(removed_bg)
        # Alpha mask
        alpha = removed_np[:, :, 3]
        # Create binary mask
        mask = cv2.threshold(alpha, 0, 255, cv2.THRESH_BINARY)[1]
        # Smooth edges
        mask = cv2.GaussianBlur(mask, (21, 21), 11)
        # Blur full image strongly
        blurred_bg = cv2.GaussianBlur(original_img, (61, 61), 30)
        # Normalize mask
        mask = mask.astype(float) / 255.0
        mask = cv2.merge([mask, mask, mask])
        # Blend foreground + blurred background
        foreground = original_img.astype(float) * mask
        background = blurred_bg.astype(float) * (1 - mask)
        final = cv2.add(foreground, background)
        final = np.uint8(final)
        # Convert final image to base64
        _, buffer = cv2.imencode('.jpg', final)
        image_base64 = base64.b64encode(buffer).decode('utf-8')
        blurred_image = image_base64
    return render(request, 'imagetool/background_blur.html', {'blurred_image': blurred_image})
##############################meme generator###################################
def MemeGenerator(request):
    meme_image = None
    if request.method == 'POST':
        uploaded_image = request.FILES.get('image')
        top_text = request.POST.get('top_text', '')
        bottom_text = request.POST.get('bottom_text', '')
        if uploaded_image:
            # Open image
            image = Image.open(uploaded_image).convert('RGB')
            draw = ImageDraw.Draw(image)
            # Font settings
            try:
                font = ImageFont.truetype("arial.ttf", 50)
            except:
                font = ImageFont.load_default()
            width, height = image.size
            # Draw Top Text
            top_position = (width // 2, 50)
            # Draw Bottom Text
            bottom_position = (width // 2, height - 80)
            # Function to draw outlined text
            def draw_text_with_outline(position, text):
                x, y = position
                # Outline
                for adj in range(-2, 3):
                    for adj2 in range(-2, 3):
                        draw.text((x + adj, y + adj2),text,font=font,fill='black',anchor='mm')
                # Main text
                draw.text((x, y),text,font=font,fill='white',anchor='mm')
            draw_text_with_outline(top_position, top_text.upper())
            draw_text_with_outline(bottom_position, bottom_text.upper())
            # Convert image to base64
            buffer = BytesIO()
            image.save(buffer, format='PNG')
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            meme_image = f"data:image/png;base64,{image_base64}"
    return render(request, 'imagetool/meme.html', {'meme_image': meme_image})

##############################pdf tool##########################################
#convert image to pdf
def convert_image_to_pdf(request):
    if request.method == 'POST' and request.FILES['image']:
        image_file = request.FILES['image']
        # Save the image temporarily
        image_path = os.path.join('media', 'temp_image.png')
        with open(image_path, 'wb') as destination:
            for chunk in image_file.chunks():
                destination.write(chunk)
        # Create a PDF file using ReportLab
        buffer = BytesIO()
        pdf = canvas.Canvas(buffer)
        # Get image dimensions
        image = Image.open(image_path)
        img_width, img_height = image.size
        # Assuming A4 paper size (you can adjust as needed)
        max_width = 600  # Adjust as needed
        max_height = 700  # Adjust as needed
        # Calculate scaling factor
        width_scale = max_width / img_width
        height_scale = max_height / img_height
        scale = min(width_scale, height_scale)
        # Calculate new dimensions
        new_width = int(img_width * scale)
        new_height = int(img_height * scale)
        # Set margins (adjust as needed)
        top_margin = 50
        bottom_margin = 50
        # Calculate centering position with margins
        x = (max_width - new_width) / 2
        y = bottom_margin  # Adjust to set the margin from the bottom
        # Draw image on the PDF with correct scaling, centering, and margins
        pdf.drawInlineImage(os.path.abspath(image_path), x, y, width=new_width, height=new_height)
        pdf.showPage()
        pdf.save()
        # Close the Image object
        image.close()
        # Remove the temporarily saved image
        os.remove(image_path)
        # Set response headers for PDF download
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="output.pdf"'
        buffer.seek(0)
        response.write(buffer.read())
        buffer.close()
        return response
    return render(request, 'pdftool/convertpdf.html')
#combine multiple image pdf
def combine_images_to_pdf(request):
    if request.method == 'POST':
        # Assuming you have a form with multiple file inputs named 'images'
        images = request.FILES.getlist('images')
        if images:
            # Create a BytesIO buffer to store the PDF content
            pdf_buffer = BytesIO()
            # Create a PDF document using reportlab
            pdf = canvas.Canvas(pdf_buffer)
            for image in images:
                # Open each image using PIL
                img = Image.open(image)
                # Convert the image to RGB mode
                img = img.convert('RGB')
                # Get image size in points (1 inch = 72 points)
                width, height = img.size
                aspect_ratio = height / float(width)
                # Set the image size in the PDF document
                pdf.setPageSize((400, 400 * aspect_ratio))
                # Draw the image on the PDF
                pdf.drawInlineImage(img, 0, 0, width=400, height=400 * aspect_ratio)
                # Add a new page for the next image
                pdf.showPage()
            # Close the PDF
            pdf.save()
            # Reset the buffer position to the beginning
            pdf_buffer.seek(0)
            # Create a response with the PDF content
            response = HttpResponse(pdf_buffer.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="combined_images.pdf"'
            return response
    return render(request, 'pdftool/combine_images.html')
#two pdf merger
def merge_pdfs(pdf1, pdf2):
    merger = PdfMerger()
    merger.append(pdf1)
    merger.append(pdf2)
    merged_pdf_path = 'media/merged_pdf.pdf'  # Adjust the path as needed
    merger.write(merged_pdf_path)
    merger.close()
    return merged_pdf_path
def merge_pdfs_view(request):
    if request.method == 'POST':
        pdf1 = request.FILES['pdf1']
        pdf2 = request.FILES['pdf2']
        merged_pdf_path = merge_pdfs(pdf1, pdf2)
        with open(merged_pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="merged_pdf.pdf"'
            return response
    return render(request, 'pdftool/mergepdf.html')
#convert text file to pdf
def convert_text_to_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    return pdf.output(dest='S').encode('latin1')
def convert_text_file_to_pdf(file):
    text_content = file.read().decode('utf-8')
    return convert_text_to_pdf(text_content)
def text_to_pdf(request):
    if request.method == 'POST':
        if 'text_file' in request.FILES:
            text_file = request.FILES['text_file']
            pdf_content = convert_text_file_to_pdf(text_file)
            response = HttpResponse(content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{text_file.name}.pdf"'
            response.write(pdf_content)
            return response
    return render(request, 'pdftool/text-pdf.html')
#convert excel file to pdf
def convert_excel_to_pdf(request):
    if request.method == 'POST' and request.FILES.get('excel_file'):
        # Handle file upload
        excel_file = request.FILES['excel_file']
        # Create the 'temp' directory if it doesn't exist
        temp_dir = os.path.join(settings.MEDIA_ROOT, 'uploadexcel')
        os.makedirs(temp_dir, exist_ok=True)
        # Save the Excel file to a temporary location
        excel_file_path = os.path.join(temp_dir, excel_file.name)
        with open(excel_file_path, 'wb') as destination:
            for chunk in excel_file.chunks():
                destination.write(chunk)
        # Convert Excel to PDF using LibreOffice
        pdf_file_path = os.path.join(temp_dir, f"{os.path.splitext(excel_file.name)[0]}.pdf")
        convert_excel_to_pdf_libreoffice(excel_file_path, pdf_file_path)
        os.remove(excel_file_path)
        # Download the PDF file
        with open(pdf_file_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.basename(pdf_file_path)}"'
            # Remove the PDF file after serving for download
            pdf_file.close()
            os.remove(pdf_file_path)
            return response
    return render(request, 'pdftool/excel-pdf.html')

def convert_excel_to_pdf_libreoffice(input_path, output_path):
    # Convert Excel to PDF using LibreOffice
    command = [
        settings.LIBREOFFICE_EXECUTABLE,
        '--headless',
        '--convert-to',
        'pdf',
        '--outdir',
        os.path.dirname(output_path),
        input_path,
    ]
    process = Popen(command, stdout=PIPE, stderr=PIPE)
    stdout, stderr = process.communicate()
    if process.returncode != 0:
        raise Exception(f"LibreOffice conversion failed: {stderr.decode('utf-8')}")
#convert csv to pdf
def convert_csv_to_pdf(request):
    if request.method == 'POST' and request.FILES.get('csv_file'):
        # Handle file upload
        csv_file = request.FILES['csv_file']
        # Save the CSV file to a temporary location
        csv_file_path = settings.MEDIA_ROOT + '/uploadcsv/' + csv_file.name
        with open(csv_file_path, 'wb') as destination:
            for chunk in csv_file.chunks():
                destination.write(chunk)
        # Convert CSV to PDF using LibreOffice
        pdf_file_path = csv_file_path.rsplit('.', 1)[0] + '.pdf'
        convert_csv_to_pdf_libreoffice(csv_file_path, pdf_file_path)
        os.remove(csv_file_path)
        # Download the PDF file
        with open(pdf_file_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{csv_file.name.rsplit(".", 1)[0]}.pdf"'
            pdf_file.close()
            os.remove(pdf_file_path)
            return response
    return render(request, 'pdftool/csv_to_pdf.html')
def convert_csv_to_pdf_libreoffice(input_path, output_path):
    # Convert CSV to PDF using LibreOffice
    command = [
        settings.LIBREOFFICE_EXECUTABLE,
        '--headless',
        '--convert-to',
        'pdf',
        '--outdir',
        settings.MEDIA_ROOT + '/uploadcsv/',
        input_path,
    ]
    process = Popen(command, stdout=PIPE, stderr=PIPE)
    stdout, stderr = process.communicate()
    if process.returncode != 0:
        raise Exception(f"LibreOffice conversion failed: {stderr.decode('utf-8')}")
# Add the  directory 
temp_dir = settings.MEDIA_ROOT + '/uploadcsv/'
if not os.path.exists(temp_dir):
    os.makedirs(temp_dir)
#convert word file in to pdf
def convert_word_to_pdf(input_path, output_path):
    # Convert Word to PDF using LibreOffice
    command = [
        settings.LIBREOFFICE_EXECUTABLE,
        '--headless',
        '--convert-to',
        'pdf',
        '--outdir',
        os.path.dirname(output_path),
        input_path,
    ]
    process = Popen(command, stdout=PIPE, stderr=PIPE)
    stdout, stderr = process.communicate()
    if process.returncode != 0:
        raise Exception(f"LibreOffice conversion failed: {stderr.decode('utf-8')}")
def convert_word_to_pdf_and_download(request):
    if request.method == 'POST' and request.FILES.get('word_file'):
        # Handle file upload
        word_file = request.FILES['word_file']
        # Ensure the 'temp' directory exists
        temp_dir = os.path.join(settings.MEDIA_ROOT, 'uploaddoc')
        os.makedirs(temp_dir, exist_ok=True)
        # Save the Word file to a temporary location
        word_file_path = os.path.join(temp_dir, word_file.name)
        with open(word_file_path, 'wb') as destination:
            for chunk in word_file.chunks():
                destination.write(chunk)
        # Convert Word to PDF
        pdf_file_path = os.path.join(temp_dir, f"{os.path.splitext(word_file.name)[0]}.pdf")
        convert_word_to_pdf(word_file_path, pdf_file_path)
        os.remove(word_file_path)
        # Download the PDF file
        with open(pdf_file_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.basename(pdf_file_path)}"'
            pdf_file.close()
            os.remove(pdf_file_path)
            return response
    return render(request, 'pdftool/word_to_pdf.html')
#convert ppt file to pdf
def convert_ppt_to_pdf(file_path):
    pdf_file_path = os.path.splitext(file_path)[0] + '.pdf'
    # Provide the full path to 'soffice' executable
    libreoffice_path = settings.LIBREOFFICE_EXECUTABLE
    #r'C:\Program Files\LibreOffice\program\soffice.exe'  # Update this path
    # Use the full path to 'soffice' for conversion
    command = [libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(pdf_file_path), file_path]
    subprocess.run(command, shell=True)
    return pdf_file_path
def download_file(request, file_path):
    with open(file_path, 'rb') as file:
        response = HttpResponse(file.read(), content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename={os.path.basename(file_path)}'
        file.close()
        os.remove(file_path)
        return response
def convert_ppt(request):
    if request.method == 'POST' and request.FILES.get('ppt_file'):
        ppt_file = request.FILES['ppt_file']
        # Create the 'media/uploads' directory if it doesn't exist
        uploads_directory = os.path.join('media', 'uploadsppt')
        os.makedirs(uploads_directory, exist_ok=True)
        ppt_file_path = os.path.join(uploads_directory, ppt_file.name)
        with open(ppt_file_path, 'wb') as destination:
            for chunk in ppt_file.chunks():
                destination.write(chunk)
        try:
            pdf_file_path = convert_ppt_to_pdf(ppt_file_path)
            # Clean up the PPT file after conversion
            os.remove(ppt_file_path)
            return download_file(request, pdf_file_path)
        except Exception as e:
            # Log the exception to help diagnose the issue
            print(f"Conversion error: {e}")
    return render(request, 'pdftool/ppt_to_pdf.html')
#convert pdf to word file
def convert_pdf_to_word(pdf_path, output_path):
    # Using pdf2docx to convert PDF to DOCX
    cv = Converter(pdf_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
def pdf_to_word(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        pdf_path = os.path.join(settings.MEDIA_ROOT, 'input.pdf')
        output_path = os.path.join(settings.MEDIA_ROOT, 'output', 'output.docx')
        # Create the output directory if it doesn't exist
        os.makedirs(os.path.join(settings.MEDIA_ROOT, 'output'), exist_ok=True)
        with open(pdf_path, 'wb') as destination:
            for chunk in pdf_file.chunks():
                destination.write(chunk)
        convert_pdf_to_word(pdf_path, output_path)
        os.remove(pdf_path)
        # Check if the Word file exists before attempting to open it
        if os.path.exists(output_path):
            with open(output_path, 'rb') as word_file:
                response = HttpResponse(word_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                response['Content-Disposition'] = 'attachment; filename=output.docx'
                word_file.close()
                os.remove(output_path)
            return response
        else:
            return HttpResponse("Conversion failed. The Word file is not available.", status=500)
    return render(request, 'pdftool/pdf_to_word.html')
#convert pdf to ppt
def convert_pdf_to_ppt(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        # Get the uploaded PDF file from the form
        pdf_file = request.FILES['pdf_file']
        # Read the PDF file and create a PowerPoint presentation
        pdf_reader = PdfReader(pdf_file)
        presentation = Presentation()
        for page_num in range(len(pdf_reader.pages)):
            # Create a blank slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Choose a slide layout
            # Extract text from the PDF page and add it to the slide
            text = pdf_reader.pages[page_num].extract_text()
            title_shape = slide.shapes.title
            title_shape.text = text
        # Save the PowerPoint presentation to a BytesIO buffer
        ppt_buffer = BytesIO()
        presentation.save(ppt_buffer)
        ppt_buffer.seek(0)
        # Set response headers for file download
        response = HttpResponse(ppt_buffer.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename={pdf_file.name}.pptx'
        return response
    return render(request, 'pdftool/pdf_to_ppt.html')
#convert pdf to excel 
def convert_pdf_to_excel(request):
    if request.method == 'POST':
        pdf_file = request.FILES['pdf_file']
        pdf_file_path = os.path.join('media', pdf_file.name)
        with open(pdf_file_path, 'wb') as destination:
            for chunk in pdf_file.chunks():
                destination.write(chunk)
        # Use PyMuPDF to extract text from the PDF
        pdf_document = fitz.open(pdf_file_path)
        text = ""
        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            text += page.get_text()
        # Close the PyMuPDF document
        pdf_document.close()
        # Convert the extracted text to a list of rows
        data = [line.split('\t') for line in text.split('\n')]
        # Remove empty rows
        data = [row for row in data if row]
        # Determine the maximum number of columns in any row
        max_columns = max(len(row) for row in data)
        # Fill missing columns with empty strings
        data = [row + [''] * (max_columns - len(row)) for row in data]
        # Convert the list of rows to a Pandas DataFrame
        df = pd.DataFrame(data)
        # Convert DataFrame to Excel
        excel_file_path = os.path.join('media', 'output_excel_file.xlsx')
        df.to_excel(excel_file_path, index=False)
        # Remove the PDF file
        os.remove(pdf_file_path)
        # Provide the Excel file for download
        with open(excel_file_path, 'rb') as excel_file:
            response = HttpResponse(excel_file.read(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = 'attachment; filename=output_excel_file.xlsx'
        # Remove the Excel file after download
        os.remove(excel_file_path)
        return response
    return render(request, 'pdftool/pdf_to_excel.html')
#convert pdf to csv 
def convert_pdf_to_csv(pdf_path, csv_path):
    # Extract text from PDF using PyMuPDF
    pdf_document = fitz.open(pdf_path)
    text = ""
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text += page.get_text()
    # Convert text data to CSV
    with open(csv_path, 'w', newline='', encoding='utf-8') as csvfile:
        csv_writer = csv.writer(csvfile)
        # Split text into potential rows
        rows = [row.strip() for row in text.split('\n') if row.strip()]
        # Determine dynamic header and write it to CSV
        max_fields = max(len(row.split(',')) for row in rows)
        header = [f'Column_{i}' for i in range(1, max_fields + 1)]
        csv_writer.writerow(header)
        # Write data rows to CSV
        for row in rows:
            try:
                csv_writer.writerow([cell.strip() for cell in row.split(',')])
            except Exception as e:
                print(f"Error processing line: {row}")
                print(f"Error details: {e}")
def pdf_to_csv(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        # Handle the uploaded PDF file
        pdf_file = request.FILES['pdf_file']
        # Set up media directories
        media_root = 'media'
        upload_dir = os.path.join(media_root, 'uploads')
        download_dir = os.path.join(media_root, 'downloads')
        # Ensure the necessary directories exist
        os.makedirs(upload_dir, exist_ok=True)
        os.makedirs(download_dir, exist_ok=True)
        # Construct file paths
        pdf_path = os.path.join(upload_dir, pdf_file.name)
        csv_path = os.path.join(download_dir, f'{os.path.splitext(pdf_file.name)[0]}.csv')
        # Save the uploaded PDF file
        with open(pdf_path, 'wb+') as destination:
            for chunk in pdf_file.chunks():
                destination.write(chunk)
        # Convert PDF to CSV
        convert_pdf_to_csv(pdf_path, csv_path)
        # Delete PDF file after conversion
        os.remove(pdf_path)
        # Display CSV data on HTML page
        csv_data = pd.read_csv(csv_path)
        csv_html = csv_data.to_html()
        # Provide download link
        download_link = f'<a href="/download/{os.path.basename(csv_path)}">Download CSV</a>'
        # Render HTML with CSV data and download link
        return render(request, 'pdftool/pdf_to_csv.html', {'csv_html': csv_html, 'download_link': download_link})
    return render(request, 'pdftool/pdf_to_csv.html')
class AudioFileWrapper:
    def __init__(self, file_path):
        self.file_path = file_path
    def __iter__(self):
        with open(self.file_path, 'rb') as audio_file:
            for chunk in audio_file:
                yield chunk
        # Remove the csv  file after the file is closed
        os.remove(self.file_path)
def download_csv(request, file_name):
    # Set up media directory
    download_dir = os.path.join('media', 'downloads')
    file_path = os.path.join(download_dir, file_name)
    # Serve the file for download
    response = FileResponse(open(file_path, 'rb'))
    response['Content-Disposition'] = f'attachment; filename="{file_name}"'
    return response
###################################Other tool#################################################
import whois
from countryinfo import CountryInfo
from moviepy.editor import VideoFileClip
from app.forms import VideoForm
from tempfile import NamedTemporaryFile
from django.http import FileResponse
from django.views.decorators.http import require_GET
import pyshorteners
#Domain name information
def get_domain_info(request):
    domain_info = None
    if request.method == 'POST':
        domain_name = request.POST.get('domain_name', '')
        if domain_name:
            try:
                domain_info = whois.whois(domain_name)
            except whois.parser.PywhoisError as e:
                error_message = f"Domain name not register: {e}"
                return render(request, 'othertool/domain_info.html', {'error_message': error_message})
    return render(request, 'othertool/domain_info.html', {'domain_info': domain_info})
#Country information
def get_country_info(request):
    if request.method == 'POST':
        country_code = request.POST.get('country_code')
        country_info = get_country_information(country_code)
        return render(request, 'othertool/country_info.html', {'country_info': country_info})
    return render(request, 'othertool/country_info.html')
def get_country_information(country_code):
    try:
        country = CountryInfo(country_code)
        return {
            'name': country.name(),
            'capital': country.capital(),
            'population': country.population(),
            'area': country.area(),
            'currency': country.currencies(),
            'languages': country.languages(),
            'borders': country.borders(),
            'other_name': country.alt_spellings(),
            'calling_code': country.calling_codes(),
            'capital_length': country.capital_latlng(),
            'demonym': country.demonym(),
            'native_name': country.native_name(),
            'region': country.region(),
            'subregion': country.subregion(),
            'time_zones': country.timezones(),
            'more_info': country.wiki(),
            
        }
    except Exception as e:
        return {'error': f"Country name invalid {str(e)}"}
#convert video to audio
def convert_video(request):
    if request.method == 'POST':
        form = VideoForm(request.POST, request.FILES)
        if form.is_valid():
            video_file = request.FILES['video_file']
            audio_filename = convert_to_audio(video_file)
            return render(request, 'othertool/convert_video.html', {'audio_filename': audio_filename})
    else:
        form = VideoForm()
    return render(request, 'othertool/convert_video.html', {'form': form})
def convert_to_audio(video_file):
    # Create a temporary file to store the video content
    with NamedTemporaryFile(delete=False) as temp_video:
        for chunk in video_file.chunks():
            temp_video.write(chunk)
    video_clip = VideoFileClip(temp_video.name)
    audio_clip = video_clip.audio
    # Create a temporary file to store the audio content
    audio_filename = f"{video_file.name.split('.')[0]}.mp3"
    audio_filepath = os.path.join(settings.MEDIA_ROOT, audio_filename)
    audio_clip.write_audiofile(audio_filepath)
    # Close the video and audio clips
    video_clip.close()
    audio_clip.close()
    # Clean up temporary video file
    os.remove(temp_video.name)
    return audio_filename
class AudioFileWrapper:
    def __init__(self, file_path):
        self.file_path = file_path
    def __iter__(self):
        with open(self.file_path, 'rb') as audio_file:
            for chunk in audio_file:
                yield chunk
        # Remove the audio file after the file is closed
        os.remove(self.file_path)
@require_GET
def download_audio(request, audio_filename):
    audio_file_path = os.path.join(settings.MEDIA_ROOT, audio_filename)
    # Create a FileResponse with a custom file-like object
    response = FileResponse(AudioFileWrapper(audio_file_path), content_type='audio/mpeg')
    response['Content-Disposition'] = f'attachment; filename="{audio_filename}"'
    return response
#url shortners
# In-memory storage for URLs
url_mapping = {}
def Shortener_url(request):
    if request.method == 'POST':
        long_url = request.POST.get('long_url')
        if long_url:
            shortener = pyshorteners.Shortener()
            short_url = shortener.tinyurl.short(long_url)
            url_mapping[short_url] = long_url
    # Display sorted URLs
    sorted_urls = sorted(url_mapping.items())
    return render(request, 'othertool/url_shortner.html', {'urls': sorted_urls})
def redirect_original(request, short_url):
    long_url = url_mapping.get(short_url, '/')
    return redirect(long_url)
# IP Lookup
import requests
def IpLookup(request):
    data = None
    if request.method == "POST":
        ip = request.POST.get("ip")
        url = f"http://ip-api.com/json/{ip}"
        response = requests.get(url)
        data = response.json()
        return render(request, "othertool/iplookup.html", {"ip_addr": data})
    return render(request, "othertool/iplookup.html")
# QR Code Generator 
import qrcode
import base64
import time
def Qr_Generator(request):
    qr_code = None
    url = None
    if request.method == "POST":
        url = request.POST.get("url")
        qr = qrcode.make(url)
        buffer = BytesIO()
        qr.save(buffer, format="PNG")
        qr_code = base64.b64encode(buffer.getvalue()).decode()
    return render(request, "othertool/qr.html", {"qr_code": qr_code,"url": url})
#########################Seo Tool###########################   
from .utils import generate_og_image
import xml.etree.ElementTree as ET
import re
def meta_generator(request):
    context = {}
    if request.method == "POST":
        context = {
            "title": request.POST.get("title"),
            "description": request.POST.get("description"),
            "keywords": request.POST.get("keywords"),
            "url": request.POST.get("url"),
        }
    return render(request, "seotool/meta_generator.html", context)
def og_form(request):
    return render(request, 'seotool/og_form.html')
def preview_image(request):
    title = request.POST.get('title', 'Default Title')
    description = request.POST.get('description', 'Default Description')

    return render(request, 'seotool/og_preview.html', {
        'title': title,
        'description': description
    })
def download_image(request):
    title = request.GET.get('title', 'Default Title')
    description = request.GET.get('description', 'Default Description')
    image = generate_og_image(title, description)

    response = HttpResponse(image, content_type="image/png")
    response['Content-Disposition'] = 'attachment; filename="og-image.png"'
    return response
def RobotsGenerator(request):
    content = ""
    if request.method == "POST":
        user_agent = request.POST.get("user_agent", "*")
        disallow = request.POST.get("disallow", "")
        allow = request.POST.get("allow", "")
        crawl_delay = request.POST.get("crawl_delay", "")
        sitemap = request.POST.get("sitemap", "")

        lines = []
        lines.append(f"User-agent: {user_agent}")

        if disallow:
            for path in disallow.split(','):
                lines.append(f"Disallow: {path.strip()}")

        if allow:
            for path in allow.split(','):
                lines.append(f"Allow: {path.strip()}")

        if crawl_delay:
            lines.append(f"Crawl-delay: {crawl_delay}")

        if sitemap:
            lines.append(f"Sitemap: {sitemap}")

        content = "\n".join(lines)
    return render(request, "seotool/robots_text.html", {
        "robots_content": content
    })
def SitemapValidator(request):
    result = None
    urls = []
    errors = []

    if request.method == "POST":
        sitemap_url = request.POST.get("sitemap_url")
        if sitemap_url:
            try:
                response = requests.get(sitemap_url, timeout=10)
                if response.status_code != 200:
                    errors.append(f"Failed to fetch sitemap. Status: {response.status_code}")
                else:
                    try:
                        root = ET.fromstring(response.content)
                        namespace = {'ns': 'http://www.sitemaps.org/schemas/sitemap/0.9'}
                        for url_tag in root.findall('ns:url', namespace):
                            loc = url_tag.find('ns:loc', namespace)
                            if loc is not None:
                                urls.append(loc.text)
                        if urls:
                            result = "✅ Valid Sitemap"
                        else:
                            result = "⚠️ No URLs found"
                    except ET.ParseError:
                        errors.append("Invalid XML format")
            except Exception as e:
                errors.append(str(e))
        else:
            errors.append("Please enter a sitemap URL")
    return render(request, "seotool/sitemap.html", {"result": result,"urls": urls,"errors": errors})
def KeywordDensity(request):
    result = None
    if request.method == "POST":
        form = KeywordForm(request.POST)
        if form.is_valid():
            text = form.cleaned_data['text'].lower()
            keyword = form.cleaned_data['keyword'].lower()
            # Clean text
            words = re.findall(r'\b\w+\b', text)
            total_words = len(words)
            # Count keyword occurrences
            keyword_count = text.count(keyword)
            # Calculate density
            density = (keyword_count / total_words) * 100 if total_words > 0 else 0
            result = {
                'total_words': total_words,
                'keyword_count': keyword_count,
                'density': round(density, 2),
                'keyword': keyword
            }
    else:
        form = KeywordForm()
    return render(request, 'seotool/keyword_density.html', {'form': form,'result': result})
def PageSpeedChecker(request):
    result = None
    if request.method == "POST":
        form = URLForm(request.POST)
        if form.is_valid():
            url = form.cleaned_data['url']
            try:
                start_time = time.time()
                response = requests.get(url, timeout=10)
                load_time = time.time() - start_time
                page_size = len(response.content) / 1024  # KB
                result = {
                    'url': url,
                    'status_code': response.status_code,
                    'load_time': round(load_time, 3),
                    'page_size': round(page_size, 2),
                }
            except requests.exceptions.RequestException as e:
                result = {
                    'error': str(e)
                }
    else:
        form = URLForm()
    return render(request, 'seotool/speed_checker.html', { 'form': form,'result': result})
    